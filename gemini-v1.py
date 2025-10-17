from flask import Flask, request, jsonify
import os
import PyPDF2
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import re
import google.generativeai as genai

app = Flask(__name__)

# Configure Gemini - REPLACE WITH YOUR API KEY
GEMINI_API_KEY = ""
genai.configure(api_key=GEMINI_API_KEY)

class ConciseStudyAssistant:
    def __init__(self):
        print("🎯 Starting Concise Study Assistant - Short & Focused Answers!")
        
        # Load embedding model for document search
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Initialize Gemini model
        try:
            self.gemini_model = genai.GenerativeModel('gemini-2.0-flash-exp')
            print("✅ Gemini model loaded successfully!")
        except Exception as e:
            print(f"❌ Error loading Gemini: {e}")
            raise
        
        # Initialize storage
        self.documents = []
        self.full_text = ""
        self.index = None
        self.current_file = None
        
    def extract_text(self, file_path, filename):
        """Extract text from files"""
        text = ""
        
        try:
            if filename.endswith('.pdf'):
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                        
            elif filename.endswith('.docx'):
                doc = Document(file_path)
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                    
            elif filename.endswith('.pptx'):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
            print(f"📄 Extracted {len(text)} characters from {filename}")
            self.full_text = text
            return text
            
        except Exception as e:
            print(f"❌ Error: {e}")
            return ""
    
    def chunk_text(self, text):
        """Split text into meaningful chunks"""
        paragraphs = text.split('\n\n')
        chunks = []
        
        for para in paragraphs:
            if len(para.strip()) > 50:
                chunks.append(para.strip())
        
        return chunks
    
    def process_file(self, file_path, filename):
        """Process uploaded file"""
        text = self.extract_text(file_path, filename)
        if not text:
            return False
        
        chunks = self.chunk_text(text)
        self.documents = chunks
        self.current_file = filename
        
        # Create embeddings for search
        embeddings = self.embedding_model.encode(chunks)
        
        # Create FAISS index
        dimension = embeddings.shape[1]
        self.index = faiss.IndexFlatIP(dimension)
        faiss.normalize_L2(embeddings)
        self.index.add(embeddings.astype(np.float32))
        
        print(f"✅ Processed {len(chunks)} chunks from {filename}")
        return True
    
    def find_relevant_text(self, question, top_k=3):
        """Find most relevant text for the question"""
        if not self.index:
            return []
        
        question_embedding = self.embedding_model.encode([question])
        faiss.normalize_L2(question_embedding)
        
        scores, indices = self.index.search(question_embedding.astype(np.float32), top_k)
        
        results = []
        for idx, score in zip(indices[0], scores[0]):
            if idx < len(self.documents) and score > 0.1:
                results.append(self.documents[idx])
        
        return results
    
    def generate_concise_answer(self, question, relevant_texts):
        """Generate SHORT, focused answers"""
        
        if not relevant_texts:
            # Very short response for no content
            return "Not covered in your notes."
        
        # Combine relevant texts
        context = "\n".join(relevant_texts[:2])
        
        # STRICT prompt for concise answers
        prompt = f"""Provide a SHORT answer (2-3 sentences MAX). Use notes first, supplement briefly if needed.

Question: {question}

Notes: {context}

Answer in 2-3 concise sentences:"""
        
        try:
            response = self.gemini_model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.1,
                    "max_output_tokens": 150,  # SHORT responses
                }
            )
            
            answer = response.text.strip()
            
            # Ensure it's short - take first 2 sentences if too long
            sentences = re.split(r'[.!?]+', answer)
            if len(sentences) > 2:
                answer = '. '.join(sentences[:2]) + '.'
            
            return answer
            
        except Exception as e:
            # Fallback: just use the first relevant chunk
            return relevant_texts[0] if relevant_texts else "Not in notes."
    
    def explain_concept_concise(self, concept):
        """Explain a concept in 3-4 sentences MAX"""
        relevant_texts = self.find_relevant_text(concept)
        
        context = "\n".join(relevant_texts[:2]) if relevant_texts else "No specific notes content."
        
        prompt = f"""Explain this concept in 3-4 SHORT sentences MAX. Be direct and focused.

Concept: {concept}

Notes Context: {context}

Concise Explanation:"""
        
        try:
            response = self.gemini_model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.1,
                    "max_output_tokens": 200,
                }
            )
            
            explanation = response.text.strip()
            
            # Ensure brevity - take first 3-4 sentences
            sentences = re.split(r'[.!?]+', explanation)
            if len(sentences) > 4:
                explanation = '. '.join(sentences[:3]) + '.'
            
            return explanation
            
        except Exception as e:
            if relevant_texts:
                return f"From notes: {relevant_texts[0]}"
            else:
                return "Concept not covered in notes."
    
    def qa_short(self, question):
        """Question-Answer in 1-2 sentences"""
        relevant_texts = self.find_relevant_text(question)
        
        context = "\n".join(relevant_texts[:1]) if relevant_texts else "No notes."
        
        prompt = f"""Answer in 1-2 SHORT sentences MAX.

Q: {question}

Notes: {context}

A:"""
        
        try:
            response = self.gemini_model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.1,
                    "max_output_tokens": 100,
                }
            )
            return response.text.strip()
        except Exception as e:
            return relevant_texts[0] if relevant_texts else "No info."

# Initialize chatbot
chatbot = ConciseStudyAssistant()

@app.route('/upload', methods=['POST'])
def upload_file():
    """Upload a file"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        file_path = f"temp_{file.filename}"
        file.save(file_path)
        
        success = chatbot.process_file(file_path, file.filename)
        os.remove(file_path)
        
        if not success:
            return jsonify({'error': 'Could not process file'}), 400
        
        return jsonify({
            'success': True,
            'message': 'File uploaded! I will give SHORT, focused answers.',
            'filename': file.filename,
            'chunks': len(chatbot.documents)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    """Ask questions - get SHORT answers"""
    try:
        data = request.json
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Please provide a question'}), 400
        
        if not chatbot.documents:
            return jsonify({'error': 'Please upload a document first'}), 400
        
        # Find relevant text
        relevant_texts = chatbot.find_relevant_text(question)
        
        # Generate CONCISE answer
        answer = chatbot.generate_concise_answer(question, relevant_texts)
        
        return jsonify({
            'question': question,
            'answer': answer,
            'length': len(answer),
            'chunks_used': len(relevant_texts)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/explain', methods=['POST'])
def explain_concept():
    """Explain a concept - SHORT version"""
    try:
        data = request.json
        concept = data.get('concept', '').strip()
        
        if not concept:
            return jsonify({'error': 'Please provide a concept'}), 400
        
        if not chatbot.documents:
            return jsonify({'error': 'Please upload a document first'}), 400
        
        explanation = chatbot.explain_concept_concise(concept)
        
        return jsonify({
            'concept': concept,
            'explanation': explanation,
            'length': len(explanation)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/qa', methods=['POST'])
def qa_short():
    """Very short Q&A - 1-2 sentences"""
    try:
        data = request.json
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Please provide a question'}), 400
        
        if not chatbot.documents:
            return jsonify({'error': 'Please upload a document first'}), 400
        
        answer = chatbot.qa_short(question)
        
        return jsonify({
            'question': question,
            'answer': answer,
            'length': len(answer),
            'type': 'short_qa'
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({
        'status': 'ready',
        'has_document': chatbot.current_file is not None,
        'current_file': chatbot.current_file,
        'style': 'CONCISE answers (2-3 sentences max)'
    })

if __name__ == '__main__':
    print("🎯 CONCISE Study Assistant Ready!")
    print("📝 I give SHORT, focused answers - no long essays!")
    print("\n📚 Endpoints:")
    print("   POST /upload  - Upload study materials")
    print("   POST /ask     - Ask questions (2-3 sentence answers)")
    print("   POST /explain - Explain concepts (3-4 sentences)")
    print("   POST /qa      - Very short Q&A (1-2 sentences)")
    print("   GET  /health  - Check status")
    
    print("\n🎯 Response Examples:")
    print("   ❌ Before: 500+ words, multiple concepts")
    print("   ✅ Now: 2-3 focused sentences")
    
    app.run(debug=True, host='0.0.0.0', port=5000)