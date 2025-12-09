import uuid
import subprocess
import tempfile
import io,time
import os
import base64
import random
from datetime import datetime
import uvicorn
import boto3
import fitz
import pdfplumber
from PIL import Image
from docx import Document
from pptx import Presentation
import openpyxl
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from typing import List
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel

from groq import Groq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_postgres import PGVector
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
import psycopg2 as db
from dotenv import load_dotenv
from fastapi import Query

# ---------------- ENV & CONFIG ----------------
load_dotenv()

EMBED_MODEL_NAME = "all-MiniLM-L6-v2"
CLIP_MODEL_NAME = "clip-ViT-B-32"
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
GROQ_MODEL = os.environ.get("GROQ_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct")
TABLE_NAME = "langchain_pg_embedding"

DB_HOST = os.environ.get("DB_HOST")
DB_PORT = os.environ.get("DB_PORT")
DB_NAME = os.environ.get("DB_NAME")
DB_USER = os.environ.get("DB_USER")
DB_PASSWORD = os.environ.get("DB_PASSWORD")

S3_BUCKET_NAME = os.environ.get("S3_BUCKET_NAME")
AWS_REGION = os.environ.get("AWS_REGION")
AWS_ACCESS_KEY_ID = os.environ.get("AWS_ACCESS_KEY_ID")
AWS_SECRET_ACCESS_KEY = os.environ.get("AWS_SECRET_ACCESS_KEY")

# ---------------- FASTAPI APP ----------------
app = FastAPI(title="Study Buddy API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
# ---------- Connect to hosted PostgreSQL ----------
print("Connecting to hosted PostgreSQL...")
try:
    conn = db.connect(
        host=DB_HOST,
        port=DB_PORT,
        dbname=DB_NAME,
        user=DB_USER,
        password=DB_PASSWORD,
        sslmode="require"
    )
    print("Connected to PostgreSQL!")
except Exception as e:
    print("Failed to connect:", e)
    raise

# Load embedding models
text_embeddings = HuggingFaceEmbeddings(model_name=EMBED_MODEL_NAME)
clip_model = SentenceTransformer('sentence-transformers/clip-ViT-B-32')

EMBED_DIM = len(text_embeddings.embed_query("test"))
CLIP_EMBED_DIM = 512
print(f"Text embedding dimension: {EMBED_DIM}")
print(f"CLIP embedding dimension: {CLIP_EMBED_DIM}")

# ---------------- S3 ----------------
s3 = boto3.client(
    "s3",
    aws_access_key_id=AWS_ACCESS_KEY_ID,
    aws_secret_access_key=AWS_SECRET_ACCESS_KEY,
    region_name=AWS_REGION
)
# ---------- PGVector ----------
vector_store = None
clip_vector_store = None

def get_vector_store():
    global vector_store
    if vector_store is None:
        vector_store = PGVector(
            connection=f"postgresql+psycopg2://{DB_USER}:{DB_PASSWORD}@{DB_HOST}:{DB_PORT}/{DB_NAME}?sslmode=require",
            collection_name=TABLE_NAME,
            embeddings=text_embeddings,
            distance_strategy="cosine",
            use_jsonb=True,
            pre_delete_collection=False
        )
    return vector_store

# Import Office document libraries
try:
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("DOCX support disabled: install python-docx")

try:
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("PPTX support disabled: install python-pptx")

try:
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    print("XLSX support disabled: install openpyxl")

# Import PDF processing library
try:
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("PDF support disabled: install pdfplumber")


def get_clip_vector_store():
    global clip_vector_store
    if clip_vector_store is None:
        class ClipEmbeddings:
            def embed_documents(self, texts):
                return [self.embed_query(text) for text in texts]
            
            def embed_query(self, text):
                return clip_model.encode([text])[0].tolist()
        
        clip_embeddings = ClipEmbeddings()
        clip_vector_store = PGVector(
            connection=f"postgresql+psycopg2://{DB_USER}:{DB_PASSWORD}@{DB_HOST}:{DB_PORT}/{DB_NAME}?sslmode=require",
            collection_name=f"{TABLE_NAME}_clip",
            embeddings=clip_embeddings,
            distance_strategy="cosine",
            use_jsonb=True,
            pre_delete_collection=False,
            embedding_length=512
        )
    return clip_vector_store

def retrieve_by_file_ids(file_ids, query, k=8):
    store = get_vector_store()
    clip_store = get_clip_vector_store()
    
    if not file_ids:
        return [], []
    
    filter_cond = {"file_id": file_ids}
    
    print(f"\n=== RETRIEVING CONTENT FOR FLASHCARDS/MCQs ===")
    print(f"Query: '{query}'")
    print(f"File IDs: {file_ids}")
    
    # Get text chunks
    text_results = []
    try:
        text_docs = store.similarity_search(query, k=k, filter=filter_cond)
        text_results = [doc.page_content for doc in text_docs]
        print(f"Found {len(text_results)} text chunks")
    except Exception as e:
        print(f"Text search error: {e}")
    
    # Get image descriptions - MORE AGGRESSIVE SEARCH
    image_results = []
    
    # Strategy 1: Try the query directly
    try:
        image_docs = clip_store.similarity_search(query, k=min(4, k//2), filter=filter_cond)
        for doc in image_docs:
            if doc.page_content not in image_results:
                image_results.append(doc.page_content)
        print(f"Strategy 1: Found {len(image_docs)} images with direct query")
    except Exception as e:
        print(f"Strategy 1 error: {e}")
    
    # Strategy 2: If query is specific, also try broader terms
    if query and len(query.split()) > 2:
        try:
            simple_terms = query.split()[:2]  # Take first 2 words
            for term in simple_terms:
                docs = clip_store.similarity_search(term, k=2, filter=filter_cond)
                for doc in docs:
                    if doc.page_content not in image_results:
                        image_results.append(doc.page_content)
            print(f"Strategy 2: Found additional images with simple terms")
        except Exception as e:
            print(f"Strategy 2 error: {e}")
    
    # Strategy 3: Always search for "image" and "diagram"
    try:
        visual_terms = ["image", "diagram", "chart", "graph", "figure", "photo", "picture", "illustration"]
        for term in visual_terms:
            docs = clip_store.similarity_search(term, k=2, filter=filter_cond)
            for doc in docs:
                if doc.page_content not in image_results:
                    image_results.append(doc.page_content)
        print(f"Strategy 3: Found images with visual terms")
    except Exception as e:
        print(f"Strategy 3 error: {e}")
    
    # Strategy 4: If still no images, get ANY images from these files
    if len(image_results) == 0:
        try:
            # Get random content from these files
            all_docs = clip_store.similarity_search("a", k=5, filter=filter_cond)
            for doc in all_docs:
                if doc.page_content not in image_results:
                    image_results.append(doc.page_content)
            print(f"Strategy 4: Found {len(image_results)} images with generic search")
        except Exception as e:
            print(f"Strategy 4 error: {e}")
    
    print(f"Total image descriptions retrieved: {len(image_results)}")
    
    # Show first few descriptions
    if image_results:
        print("Sample image descriptions:")
        for i, desc in enumerate(image_results[:3]):
            print(f"  {i+1}. {desc[:100]}...")
    
    print(f"=== RETRIEVAL COMPLETE ===\n")
    
    return text_results, image_results

# ---------------- DOCUMENT PROCESSING FUNCTIONS ----------------

def convert_office_to_pdf(file_stream, filename):
    """Convert Office files to PDF using local LibreOffice"""
    try:
        # Save uploaded file to temp location
        with tempfile.NamedTemporaryFile(delete=False, suffix='.tmp') as temp_input:
            file_stream.seek(0)
            temp_input.write(file_stream.read())
            temp_input_path = temp_input.name
        
        # Convert using LibreOffice
        cmd = [
            'libreoffice', '--headless', '--convert-to', 'pdf', 
            '--outdir', '/tmp', temp_input_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        
        if result.returncode == 0:
            # Find the converted PDF
            pdf_path = temp_input_path.replace('.tmp', '.pdf')
            if os.path.exists(pdf_path):
                with open(pdf_path, 'rb') as pdf_file:
                    pdf_data = pdf_file.read()
                
                # Cleanup temporary files
                os.unlink(temp_input_path)
                os.unlink(pdf_path)
                
                return pdf_data, None
        
        return None, f"Conversion failed: {result.stderr}"
        
    except subprocess.TimeoutExpired:
        return None, "Conversion timeout - LibreOffice took too long"
    except Exception as e:
        return None, f"Conversion error: {str(e)}"

def extract_text_from_pdf(file_stream):
    """Extract text from PDF file"""
    text = ""
    try:
        with pdfplumber.open(file_stream) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return ""

def extract_images_from_pdf(file_stream):
    """Extract images from PDF using PyMuPDF"""
    images = []
    try:
        # Reset stream position
        file_stream.seek(0)
        doc = fitz.open(stream=file_stream.read(), filetype="pdf")
        
        print(f"  Scanning PDF with {len(doc)} pages for images...")
        total_images = 0
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images()
            
            if image_list:
                print(f"  Page {page_num + 1}: Found {len(image_list)} images")
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # RGB or similar
                        img_data = pix.tobytes("png")
                        img_pil = Image.open(io.BytesIO(img_data))
                        images.append(img_pil)
                        total_images += 1
                        print(f"    Extracted image {total_images}: {img_pil.size} pixels")
                    
                    pix = None  # Free memory
                except Exception as img_error:
                    print(f"    Failed to extract image {img_index}: {img_error}")
                    continue
        
        doc.close()
        print(f"  Total images extracted: {len(images)}")
        return images
        
    except Exception as e:
        print(f"Error extracting images from PDF: {e}")
        return []
    
def extract_text_from_docx(file_stream):
    """Extract text from Word documents"""
    try:
        doc = Document(file_stream)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    text += f"Table: {row_text}\n"
        
        return text.strip()
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pptx(file_stream):
    """Extract text from PowerPoint presentations including speaker notes"""
    try:
        prs = Presentation(file_stream)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = []
            slide_notes = []
            
            # Extract text from slide shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            # Extract text from notes (speaker notes)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_notes.append(notes_text)
            
            # Extract text from tables
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            if cell.text:
                                row_texts.append(cell.text.strip())
                        if row_texts:
                            table_rows.append(" | ".join(row_texts))
                    if table_rows:
                        slide_content.append(f"Table: {'; '.join(table_rows)}")
            
            # Build slide text
            if slide_content or slide_notes:
                slide_text = f"=== SLIDE {slide_num + 1} ===\n"
                if slide_content:
                    slide_text += "Content:\n" + "\n".join(f"- {item}" for item in slide_content) + "\n"
                if slide_notes:
                    slide_text += "Notes:\n" + "\n".join(f"- {note}" for note in slide_notes) + "\n"
                text_parts.append(slide_text)
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        # Try a simpler approach as fallback
        try:
            prs = Presentation(file_stream)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                # Try direct notes access
                try:
                    if slide.notes_slide:
                        text += "NOTES: " + slide.notes_slide.notes_text_frame.text + "\n"
                except:
                    pass
            return text.strip()
        except:
            return ""
def extract_text_from_xlsx(file_stream):
    """Extract text from Excel files"""
    try:
        wb = openpyxl.load_workbook(file_stream)
        text = ""
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"--- Sheet: {sheet_name} ---\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(row_data)
                if row_text.strip():
                    text += row_text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"Error extracting XLSX text: {e}")
        return ""

def generate_image_description(image):
    """Generate text description of image using Groq for CLIP indexing"""
    client = Groq(api_key=GROQ_API_KEY)
    
    try:
        # Convert image to PNG bytes
        img_byte_arr = io.BytesIO()
        
        # Handle different image modes
        if image.mode in ('RGBA', 'LA', 'P'):
            # Convert to RGB for PDF compatibility
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'RGBA':
                background.paste(image, mask=image.split()[-1])
            else:
                background.paste(image)
            image = background
        elif image.mode != 'RGB':
            image = image.convert("RGB")
        
        image.save(img_byte_arr, format='PNG', optimize=True)
        img_b64 = base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
        
        print(f"  Generating description for image ({image.size[0]}x{image.size[1]}, mode: {image.mode})...")
        
        # Try different approaches if the first fails
        prompts = [
            "Describe this image in detail for search purposes. Focus on content, objects, text, and visual elements:",
            "What does this image show? Describe it thoroughly for document searching:",
            "Analyze this image and create a detailed description suitable for search indexing:"
        ]
        
        for attempt, prompt in enumerate(prompts):
            try:
                response = client.chat.completions.create(
                    model=GROQ_MODEL,
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}}
                        ]
                    }],
                    max_tokens=300,  # Increased for better descriptions
                    temperature=0.3,
                    timeout=30  # Add timeout
                )
                
                description = response.choices[0].message.content.strip()
                
                if description and len(description) > 20:  # Valid description
                    print(f"  Generated description ({len(description)} chars): {description[:100]}...")
                    return description
                else:
                    print(f"  Attempt {attempt+1}: Description too short: '{description}'")
                    
            except Exception as e:
                print(f"  Attempt {attempt+1} failed: {e}")
                continue
        
        # If all attempts fail, create a basic description
        print(f"  All attempts failed, creating basic description")
        return f"An image of size {image.size[0]}x{image.size[1]} pixels"
        
    except Exception as e:
        print(f"  Error generating image description: {e}")
        return f"Image content (error: {str(e)[:50]})"
    
def convert_image_or_text_to_pdf(file_stream, filename):
    """
    Converts image or text files to PDF.
    """
    ext = filename.lower().split('.')[-1]

    # IMAGE to PDF
    if ext in ["png", "jpg", "jpeg", "bmp", "gif", "webp"]:
        try:
            file_stream.seek(0)
            image = Image.open(file_stream)
            
            print(f"    Converting image to PDF: {image.size}, {image.mode}")
            
            # Handle different image modes
            if image.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', image.size, (255, 255, 255))
                if image.mode == 'RGBA':
                    background.paste(image, mask=image.split()[-1])
                else:
                    background.paste(image)
                image = background
            elif image.mode != 'RGB':
                image = image.convert("RGB")
            
            pdf_out = io.BytesIO()
            # Use optimal settings
            image.save(pdf_out, format="PDF", quality=95, optimize=True)
            pdf_bytes = pdf_out.getvalue()
            print(f"    PDF created: {len(pdf_bytes)} bytes")
            return pdf_bytes, None
            
        except Exception as e:
            print(f"    Image to PDF failed: {e}")
            return None, f"Image to PDF failed: {e}"

    # TEXT to PDF
    if ext == "txt":
        try:
            file_stream.seek(0)
            text = file_stream.read().decode("utf-8", errors="ignore")
            pdf_out = io.BytesIO()
            c = canvas.Canvas(pdf_out, pagesize=letter)
            
            # Set font
            c.setFont("Helvetica", 12)
            
            y = 750
            line_height = 15
            lines = text.split("\n")
            
            for line in lines:
                # Wrap long lines
                if len(line) > 100:
                    for i in range(0, len(line), 100):
                        c.drawString(30, y, line[i:i+100])
                        y -= line_height
                        if y < 40:
                            c.showPage()
                            y = 750
                else:
                    c.drawString(30, y, line)
                    y -= line_height
                    if y < 40:
                        c.showPage()
                        y = 750
            
            c.save()
            pdf_bytes = pdf_out.getvalue()
            print(f"    Text PDF created: {len(pdf_bytes)} bytes")
            return pdf_bytes, None
            
        except Exception as e:
            print(f"    Text to PDF failed: {e}")
            return None, f"Text to PDF failed: {e}"

    print(f"    Unsupported file type for conversion: {ext}")
    return None, f"Unsupported file type for conversion to PDF: {ext}"


def process_file_content(file, filename):
    """
    Process uploaded file and return:
    - text content
    - error message if any
    - list of images (PIL.Image)
    - PDF bytes (if converted or original PDF)
    """
    file_extension = filename.lower().split('.')[-1] if '.' in filename else ''
    file_bytes = file.read()
    file_stream = io.BytesIO(file_bytes)
    
    print(f"\n📄 PROCESSING FILE: {filename}, Type: {file_extension}")
    
    # ---------- Office files (convert to PDF first) ----------
    office_types = ['docx', 'pptx', 'xlsx', 'doc', 'ppt', 'xls']
    if file_extension in office_types:
        print(f"Converting {file_extension.upper()} to PDF for full processing...")
        pdf_data, error = convert_office_to_pdf(file_stream, filename)
        if error:
            print(f"Office to PDF failed: {error}, falling back to native extraction...")
            # Fallback text extraction
            try:
                file_stream.seek(0)
                if file_extension == 'docx' and DOCX_SUPPORT:
                    text = extract_text_from_docx(file_stream)
                    print(f"Fallback: Extracted {len(text)} chars from DOCX")
                    return text, f"Fallback text only. {error}", [], None
                elif file_extension == 'pptx' and PPTX_SUPPORT:
                    text = extract_text_from_pptx(file_stream)
                    print(f"Fallback: Extracted {len(text)} chars from PPTX")
                    return text, f"Fallback text only. {error}", [], None
                elif file_extension == 'xlsx' and XLSX_SUPPORT:
                    text = extract_text_from_xlsx(file_stream)
                    print(f"Fallback: Extracted {len(text)} chars from XLSX")
                    return text, f"Fallback text only. {error}", [], None
                else:
                    return None, f"No extraction fallback available. {error}", [], None
            except Exception as e:
                print(f"Fallback extraction failed: {e}")
                return None, f"Fallback extraction failed: {e}", [], None

        # Success: extract text and images from PDF
        print(f"Successfully converted to PDF ({len(pdf_data)} bytes)")
        pdf_stream = io.BytesIO(pdf_data)
        
        # Extract text
        pdf_stream.seek(0)
        text = extract_text_from_pdf(pdf_stream)
        print(f"Extracted {len(text)} chars of text")
        
        # Extract images
        pdf_stream.seek(0)
        images = extract_images_from_pdf(pdf_stream)
        print(f"🖼️ Extracted {len(images)} images")
        
        # Debug: Show image info
        for i, img in enumerate(images):
            print(f"  Image {i+1}: {img.size} pixels, mode: {img.mode}")
        
        return text, None, images, pdf_data

    # ---------- PDF files ----------
    elif file_extension == 'pdf':
        if not PDF_SUPPORT:
            return None, "PDF support not available. Install pdfplumber.", [], None
        
        print(f"📄 Processing PDF file directly...")
        
        # Extract text
        pdf_stream = io.BytesIO(file_bytes)
        text = extract_text_from_pdf(pdf_stream)
        print(f"Extracted {len(text)} chars of text")
        
        # Extract images
        pdf_stream.seek(0)
        images = extract_images_from_pdf(pdf_stream)
        print(f"🖼️ Extracted {len(images)} images")
        
        # Debug: Show image info
        for i, img in enumerate(images):
            print(f"  Image {i+1}: {img.size} pixels, mode: {img.mode}")
        
        return text, None, images, file_bytes

    # ---------- IMAGE files (JPG, PNG, GIF, etc.) ----------
    elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'gif', 'webp', 'tiff']:
        print(f"🖼️ Processing image file: {filename}")
        try:
            # Open and process the image
            file_stream.seek(0)
            
            # Try to open with PIL
            try:
                image = Image.open(file_stream)
                
                # Check image properties
                print(f"  Image info: size={image.size}, mode={image.mode}, format={image.format}")
                
                # Fix image if needed for consistency
                original_image = image.copy()
                
                if image.mode in ('RGBA', 'LA', 'P'):
                    print(f"  Converting {image.mode} to RGB for processing...")
                    # Convert to RGB for compatibility
                    background = Image.new('RGB', image.size, (255, 255, 255))
                    if image.mode == 'RGBA':
                        background.paste(image, mask=image.split()[-1])
                    else:
                        background.paste(image)
                    processed_image = background
                elif image.mode != 'RGB':
                    print(f"  Converting {image.mode} to RGB...")
                    processed_image = image.convert("RGB")
                else:
                    processed_image = image
                
                images = [original_image]  # Keep original for description generation
                print(f"  Image ready: size={image.size}, mode={image.mode}")
                
            except Exception as img_error:
                print(f"  Failed to open image: {img_error}")
                # Try alternative approach
                file_stream.seek(0)
                try:
                    # Try loading as bytes and creating image
                    from PIL import ImageFile
                    ImageFile.LOAD_TRUNCATED_IMAGES = True
                    image = Image.open(file_stream)
                    image.load()  # Force load
                    images = [image]
                    print(f"  Loaded with alternative method")
                except:
                    return None, f"Failed to process image: {img_error}", [], None
            
            # Convert to PDF for storage
            file_stream.seek(0)
            pdf_data, error = convert_image_or_text_to_pdf(file_stream, filename)
            if error:
                print(f"  PDF conversion failed: {error}")
                # Still return the images even if PDF conversion fails
                return "", error, images, None
            
            print(f"  Successfully generated PDF ({len(pdf_data)} bytes)")
            return "", None, images, pdf_data
            
        except Exception as e:
            print(f"Error processing image file {filename}: {e}")
            return None, f"Image processing failed: {e}", [], None

    # ---------- TEXT files ----------
    elif file_extension == 'txt':
        print(f"Processing text file: {filename}")
        try:
            file_stream.seek(0)
            text = file_stream.read().decode("utf-8", errors="ignore")
            print(f"  Read {len(text)} chars of text")
            
            # Convert to PDF for storage
            file_stream.seek(0)
            pdf_data, error = convert_image_or_text_to_pdf(file_stream, filename)
            if error:
                print(f"  PDF conversion failed: {error}")
                return text, error, [], None
            
            print(f"  Successfully generated PDF ({len(pdf_data)} bytes)")
            return text, None, [], pdf_data
            
        except Exception as e:
            print(f"Error processing text file {filename}: {e}")
            return None, f"Text processing failed: {e}", [], None

    # ---------- Plain text fallback ----------
    else:
        print(f"Unknown file type: {file_extension}, trying as text...")
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
            print(f"Processed as plain text: {len(text)} chars")
            return text, None, [], None
        except Exception as e:
            print(f"Unsupported file type: {file_extension}. {e}")
            return None, f"Unsupported file type: {file_extension}. {e}", [], None
        
# ---------- LLM (Groq API call) ----------

def groq_chat(context: str, question: str, images: list = None) -> str:
    client = Groq(api_key=GROQ_API_KEY)
    
    messages = [
        {
            "role": "system", 
            "content": "You are a helpful AI using student notes as a context you teach students. Use your knowledge/logic to explain in simple way only if required. Use the provided context to answer questions. never ever answer question unrelated to notes. Always cite your sources."
        }
    ]
    
    user_content = [{"type": "text", "text": f"CONTEXT:\n{context}\n\nQUESTION:\n{question}\n\nAnswer based on the context above. End with: SOURCE: [file(s)]"}]
    
    if images and len(images) > 0:
        for img in images:
            try:
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format='PNG')
                img_b64 = base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
                user_content.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}
                })
            except Exception as e:
                print(f"Error converting image to base64: {e}")
                continue
    
    messages.append({"role": "user", "content": user_content})
    
    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=messages,
            max_tokens=2048,
            temperature=0.1
        )
        
        content = response.choices[0].message.content
        image_info = f" with {len(images)} images" if images else ""
        print(f"Successfully got response from {GROQ_MODEL}{image_info}")
        return content
        
    except Exception as e:
        error_msg = f"Groq API error: {e}"
        print(error_msg)
        return error_msg
    
# Chat History Functions
def save_chat_history(file_id, question, answer):
    try:
        conn = get_db_connection()
        cur = conn.cursor()

        cur.execute(
            "INSERT INTO chat_history (file_id, question, answer, timestamp) VALUES (%s, %s, %s, %s)",
            (file_id, question, answer, datetime.now())
        )

        conn.commit()
        cur.close()
        conn.close()
        return True
    
    except Exception as e:
        print(f"Error saving chat history: {e}")
        try:
            conn.rollback()
        except:
            pass
        return False


def get_chat_history(file_ids, limit=10):
    try:
        conn = get_db_connection()
        cur = conn.cursor()

        placeholders = ",".join(["%s"] * len(file_ids))
        query = f"""
            SELECT file_id, question, answer, timestamp 
            FROM chat_history 
            WHERE file_id IN ({placeholders})
            ORDER BY timestamp DESC 
            LIMIT %s
        """

        cur.execute(query, file_ids + [limit])
        rows = cur.fetchall()

        cur.close()
        conn.close()

        return [{
            "file_id": r[0],
            "question": r[1],
            "answer": r[2],
            "timestamp": r[3].isoformat() if r[3] else None
        } for r in rows]

    except Exception as e:
        print(f"Error getting chat history: {e}")
        return []


def groq_chat_with_history(context: str, question: str, chat_history: list, images: list = None) -> str:
    client = Groq(api_key=GROQ_API_KEY)
    
    # Build messages with explicit conversation structure
    messages = [
        {
            "role": "system", 
            "content": "You are a helpful AI assistant that answers questions based on the provided document context. Provide clear, direct answers without any references to conversation history or whether this is the first question."
        }
    ]
    
    # Add chat history as previous messages
    for chat in chat_history:
        messages.append({"role": "user", "content": chat['question']})
        messages.append({"role": "assistant", "content": chat['answer']})
    
    # Add current context and question
    user_content = [{"type": "text", "text": f"DOCUMENT CONTEXT:\n{context}\n\nQUESTION: {question}"}]    
    if images and len(images) > 0:
        for img in images:
            try:
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format='PNG')
                img_b64 = base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
                user_content.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}
                })
            except Exception as e:
                print(f"Error converting image to base64: {e}")
                continue
    
    messages.append({"role": "user", "content": user_content})
    
    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=messages,
            max_tokens=2048,
            temperature=0.1
        )
        
        content = response.choices[0].message.content
        return content
        
    except Exception as e:
        error_msg = f"Groq API error: {e}"
        print(error_msg)
        return error_msg
    
def generate_mcq_with_groq(context: str, num_questions: int = 5):
    client = Groq(api_key=GROQ_API_KEY)
    
    prompt = f"""
    CONTENT:
    {context}
    
    Create {num_questions} multiple choice questions with 4 options each.
    
    REQUIREMENTS:
    - Questions should test understanding of key concepts
    - Each question has 4 options (A, B, C, D)
    - Only ONE correct answer per question (either A, B, C, or D)
    - Shuffle Answers like option A,B,C,D , dont always put correct answer at same place
    - Wrong answers should be plausible but incorrect
    - Cover different aspects of the content
    
    FORMAT:
    [
        {{
            "question": "clear question",
            "options": [
                "A. text",
                "B. text", 
                "C. text",
                "D. text"
            ],
            "correct_answer": "A"/"B"/"C"/"D"
        }}
    ]
    
    Return ONLY valid JSON.
    """
    
    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "Create high-quality multiple choice questions with 4 options and one correct answer."
                },
                {
                    "role": "user", 
                    "content": prompt
                }
            ],
            max_tokens=2048,
            temperature=0.7
        )
        
        content = response.choices[0].message.content.strip()
        print(f"Raw MCQ response: {content}")
        
        import json
        try:
            mcqs = json.loads(content)
            
            if isinstance(mcqs, list):
                return mcqs
            elif isinstance(mcqs, dict):
                if 'questions' in mcqs:
                    return mcqs['questions']
                else:
                    return [mcqs]
            else:
                return [{"error": "Invalid response format"}]
                
        except json.JSONDecodeError as e:
            print(f"JSON parse error: {e}")
            import re
            json_match = re.search(r'\[.*\]', content, re.DOTALL)
            if json_match:
                try:
                    return json.loads(json_match.group())
                except:
                    pass
            return [{"error": "Could not parse MCQs"}]
        
    except Exception as e:
        print(f"Groq API error: {e}")
        return [{"error": f"MCQ generation failed: {str(e)}"}]

def get_db_connection():
    return db.connect(
        host=DB_HOST,
        port=DB_PORT,
        dbname=DB_NAME,
        user=DB_USER,
        password=DB_PASSWORD,
        sslmode="require"
    )

    
#-----------Flashcard Generation with Groq-----------
def generate_flashcards_with_groq(context: str, num_flashcards: int = 5, fill_gaps: bool = False):
    """Generate flashcards - with option to fill knowledge gaps"""
    client = Groq(api_key=GROQ_API_KEY)
    
    # Choose prompt based on fill_gaps flag
    if fill_gaps:
        prompt = f"""
        CONTENT FROM USER'S NOTES:
        {context}
        
        Create {num_flashcards} SMART flashcards that fill knowledge gaps and make the student an expert.
        
        CRITICAL: Don't just repeat what's in the notes. Create cards that:
        1. Add missing information the notes don't cover
        2. Connect concepts to real-world applications
        3. Address common misunderstandings
        4. Explain WHY concepts matter
        5. Build from basic to advanced understanding
        
        For each card, include:
        - question: Thought-provoking, addresses a gap
        - answer: Comprehensive explanation that TEACHES
        - hint: Guides thinking without giving answer away
        
        Return ONLY this JSON format:
        [
            {{
                "question": "question text",
                "answer": "answer text", 
                "hint": "hint text"
            }}
        ]
        
        Example for sparse notes:
        Notes: "Force = mass × acceleration"
        Card: {{
            "question": "While your notes show F=ma, how do we calculate force when an object is on an inclined plane?",
            "answer": "On an inclined plane, we resolve weight into components. The force parallel to the plane = mg·sinθ, perpendicular = mg·cosθ. Friction force = μ·(mg·cosθ).",
            "hint": "Think about breaking gravity into components parallel and perpendicular to the surface."
        }}
        """
    else:
        prompt = f"""
        CONTENT FROM USER'S NOTES:
        {context}
        
        Create {num_flashcards} high-quality flashcards based on this content.
        
        REQUIREMENTS:
        - Questions should test understanding of key concepts
        - Answers should be clear and educational
        - Hints should guide thinking without revealing answers
        - Focus on the most important information
        
        Return ONLY this JSON format:
        [
            {{
                "question": "question text",
                "answer": "answer text", 
                "hint": "hint text"
            }}
        ]
        """
    
    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educational content creator. Create effective flashcards that help students learn."
                },
                {
                    "role": "user", 
                    "content": prompt
                }
            ],
            max_tokens=2000,
            temperature=0.7
        )
        
        content = response.choices[0].message.content.strip()
        
        import json
        try:
            flashcards = json.loads(content)
            
            if isinstance(flashcards, list):
                return flashcards
            elif isinstance(flashcards, dict):
                if 'flashcards' in flashcards:
                    return flashcards['flashcards']
                else:
                    return [flashcards]
            else:
                return [{"error": "Invalid response format"}]
                
        except json.JSONDecodeError as e:
            print(f"JSON parse error: {e}")
            import re
            json_match = re.search(r'\[.*\]', content, re.DOTALL)
            if json_match:
                try:
                    return json.loads(json_match.group())
                except:
                    pass
            return [{"error": "Could not parse flashcards"}]
        
    except Exception as e:
        print(f"Groq API error: {e}")
        return [{"error": f"Flashcard generation failed: {str(e)}"}]
    
def generate_missing_notes(context: str):
    """Generate additional notes to fill knowledge gaps"""
    client = Groq(api_key=GROQ_API_KEY)
    
    prompt = f"""
    STUDENT'S CURRENT NOTES:
    {context}
    
    TASK: Create additional notes that fill knowledge gaps and make these notes complete.
    
    INSTRUCTIONS:
    1. Analyze what's MISSING or shallow in these notes
    2. Create 3-5 sections of additional educational content
    3. Each section should address a specific gap
    4. Write in clear, study-friendly language
    5. Include examples and key takeaways
    
    IMPORTANT: Do NOT repeat what's already in the notes. Only add NEW, VALUABLE information.
    
    Return ONLY this JSON format:
    {{
        "missing_notes": [
            {{
                "title": "Section title",
                "content": "Educational content that fills a gap",
                "why_important": "Why this was missing from original notes"
            }}
        ],
        "study_advice": "Brief advice on how to study with these enhanced notes"
    }}
    
    Example response:
    {{
        "missing_notes": [
            {{
                "title": "Real-World Applications of Newton's Laws",
                "content": "Newton's laws apply everywhere: 1) Car safety (seatbelts), 2) Sports (throwing balls), 3) Space travel...",
                "why_important": "Original notes only gave formulas, no applications"
            }}
        ],
        "study_advice": "Add these sections after your formula notes to understand practical use"
    }}
    """
    
    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert tutor enhancing student notes with missing information."
                },
                {
                    "role": "user", 
                    "content": prompt
                }
            ],
            max_tokens=2500,
            temperature=0.6
        )
        
        import json
        result = json.loads(response.choices[0].message.content)
        return result
        
    except Exception as e:
        print(f"Groq API error in notes generation: {e}")
        return {
            "missing_notes": [],
            "study_advice": "Could not generate additional notes. Please try again."
        }
# ------------------- ENDPOINTS --------------------------------------------------------

@app.get("/")
def health():
    try:
        client = Groq(api_key=GROQ_API_KEY)
        models = client.models.list()
        model_names = [model.id for model in models.data]

        return {
            "status": "ok",
            "groq": "connected",
            "model": GROQ_MODEL,
            "clip_ready": True,
            "office_support": {
                "docx": DOCX_SUPPORT,
                "pptx": PPTX_SUPPORT,
                "xlsx": XLSX_SUPPORT,
                "pdf": PDF_SUPPORT,
                "libreoffice": True
            },
            "available_models": model_names
        }
    except Exception as e:
        return JSONResponse(
            content={"status": "ok", "groq": "disconnected", "error": str(e)},
            status_code=200
        )


SUPPORTED_EXTENSIONS = {
    'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls',
    'pdf', 'png', 'jpg', 'jpeg', 'bmp', 'gif', 'webp', 'tiff', 'txt'
}

def get_file_extension(filename: str) -> str:
    """Extract file extension from filename."""
    if not filename or '.' not in filename:
        return ''
    return filename.lower().split('.')[-1]

def is_file_supported(filename: str) -> bool:
    """Check if file extension is supported."""
    extension = get_file_extension(filename)
    return extension in SUPPORTED_EXTENSIONS

@app.post("/upload-files")
async def upload_files(
    files: list[UploadFile] = File(...),
    moduleId: str = Form(None),
    ocr: bool = Form(False)
):
    results = []
    errors = []

    # Validate all files first
    for file in files:
        if not is_file_supported(file.filename):
            errors.append({
                "file_name": file.filename,
                "error": f"Unsupported file type: {get_file_extension(file.filename)}. Supported formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            })

    if len(errors) == len(files):
        return JSONResponse(
            status_code=400,
            content={"error": "No supported files provided", "errors": errors}
        )

    # Get vector stores
    store = get_vector_store()
    clip_store = get_clip_vector_store()

    for file in files:
        filename = file.filename
        
        if not is_file_supported(filename):
            continue
            
        content = await file.read()
        file_stream = io.BytesIO(content)
        
        print(f"\n=== UPLOADING FILE: {filename} ===")
        
        # Extract text, images, pdf_data (YOUR ORIGINAL CODE)
        text, error, images, pdf_data = process_file_content(file_stream, filename)
        
        print(f"Processing results - Text length: {len(text) if text else 0}, Images found: {len(images)}, Error: {error}")
        
        file_id = str(uuid.uuid4())
        timestamp = int(time.time() * 1000)
        pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
        s3_key = f"modules/{moduleId}/{timestamp}-{pdf_filename}" if moduleId else f"uploads/{file_id}-{pdf_filename}"
        s3_url = f"https://{S3_BUCKET_NAME}.s3.{AWS_REGION}.amazonaws.com/{s3_key}"

        try:
            if pdf_data:
                s3.put_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=s3_key,
                    Body=pdf_data,
                    ContentType='application/pdf'
                )
                print(f"Uploaded PDF to S3: {s3_key}")
            else:
                s3.put_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=s3_key,
                    Body=content,
                    ContentType=file.content_type
                )
        except Exception as e:
            print(f"S3 upload failed: {e}")
            errors.append({"file_name": filename, "error": f"S3 upload failed: {e}"})
            s3_key = None
            s3_url = None

        chunks = []
        # OCR PROCESSING
        ocr_text = ""
        if ocr:
            print(f"OCR flag enabled, extracting text with Groq...")
            try:
                # Create FRESH stream for OCR (don't reuse the consumed stream)
                ocr_stream = io.BytesIO(content)
                ocr_text = extract_text_with_groq_ocr(ocr_stream, filename)
                print(f"OCR extracted {len(ocr_text)} characters")
                
                # Use OCR text for vector storage
                if ocr_text and len(ocr_text) > 10:
                    text = ocr_text
                    print(f"Using OCR text for vector storage")
            except Exception as e:
                print(f"OCR failed: {e}")
                ocr_text = f"OCR failed: {str(e)}"
        
        # Store text chunks in vector DB
        if text and text.strip():
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            chunks = splitter.split_text(text)

            texts = chunks
            metadatas = [{
                "file_id": file_id,
                "file_name": filename,
                "chunk_index": i,
                "file_type": filename.lower().split('.')[-1] if '.' in filename else 'unknown',
                "content_type": "text",
                "has_images": len(images) > 0,
                "ocr_processed": ocr
            } for i in range(len(chunks))]
            ids = [str(uuid.uuid4()) for _ in chunks]
            
            try:
                store.add_texts(texts=texts, metadatas=metadatas, ids=ids)
                print(f"Stored {len(chunks)} text chunks in vector DB")
            except Exception as e:
                print(f"Failed to store text chunks: {e}")

        # Store image descriptions in CLIP vector store
        if images and len(images) > 0:
            print(f"Processing {len(images)} images for CLIP storage...")
            images_stored = 0
            
            for img_index, img in enumerate(images):
                print(f"  Processing image {img_index + 1}/{len(images)}...")
                try:
                    # Generate image description
                    description = generate_image_description(img)
                    print(f"  Generated description: {description[:100]}...")
                    
                    # Prepare metadata
                    metadata = {
                        "file_id": file_id,
                        "file_name": filename,
                        "image_index": img_index,
                        "file_type": "image",
                        "content_type": "image",
                        "original_content": description[:500]  # Limit description length
                    }
                    
                    # Store in CLIP vector store
                    clip_store.add_texts(
                        texts=[description],
                        metadatas=[metadata],
                        ids=[str(uuid.uuid4())]
                    )
                    
                    images_stored += 1
                    print(f"  Successfully stored image {img_index + 1}")
                    
                except Exception as e:
                    print(f"  Failed to process/store image {img_index}: {e}")
            
            print(f"Total images stored in CLIP: {images_stored}/{len(images)}")
        else:
            print(f"No images found to store in CLIP")

        results.append({
            "file_name": filename,
            "file_id": file_id,
            "chunks": len(chunks),
            "images_processed": len(images),
            "file_type": filename.lower().split('.')[-1] if '.' in filename else 'unknown',
            "s3_key": s3_key,
            "s3_url": s3_url,
            "html": None,
            "error": error,
            "debug_info": {
                "text_extracted": len(text) if text else 0,
                "images_found": len(images),
                "images_stored_in_clip": images_stored if 'images_stored' in locals() else 0
            }
        })
        
        print(f"=== COMPLETED: {filename} ===\n")

    response_data = {"uploaded": results}
    if errors:
        response_data["errors"] = errors

    return JSONResponse(content=response_data, status_code=201)

class AskRequest(BaseModel):
    question: str
    file_ids: list[str]


@app.post("/ask")
def ask(data: AskRequest):
    question = data.question.strip()
    file_ids = data.file_ids

    if not file_ids:
        raise HTTPException(status_code=400, detail="file_ids are required")

    # 1. Get relevant content (text + images)
    text_chunks, image_descriptions = retrieve_by_file_ids(file_ids, question, k=8)
    
    # DEBUG: Print what we retrieved
    print(f"DEBUG: Retrieved {len(text_chunks)} text chunks, {len(image_descriptions)} image descriptions")
    
    if image_descriptions:
        print(f"DEBUG: First image description: {image_descriptions[0][:100]}...")
    
    # 2. Build BETTER context with clear labeling
    context_parts = []
    
    if text_chunks:
        # Add text with clear identifier
        text_context = "TEXT FROM DOCUMENTS:\n" + "\n---\n".join(text_chunks)
        context_parts.append(text_context)
    
    if image_descriptions:
        # Add image descriptions with clear identifier
        image_context = "IMAGES IN DOCUMENTS:\n"
        for i, desc in enumerate(image_descriptions):
            # Clean up the description
            clean_desc = desc.replace("This image shows", "Shows").replace("The image depicts", "Depicts")
            image_context += f"- {clean_desc}\n"
        context_parts.append(image_context)
    
    context = "\n\n".join(context_parts) if context_parts else "No content found."
    
    # 3. Get chat history
    chat_history = get_chat_history(file_ids, limit=3)
    
    # 4. Call Groq with BETTER system prompt
    client = Groq(api_key=GROQ_API_KEY)
    
    messages = [
        {
            "role": "system", 
            "content": "You are analyzing document content provided by the user. The user will provide text excerpts and descriptions of any images/diagrams. Answer their questions using this provided content. When referring to visual content, do so naturally (e.g., 'the diagram shows', 'as seen in the graph', 'the illustration demonstrates')."
        }
    ]
    
    # Add chat history
    for chat in chat_history:
        messages.append({"role": "user", "content": chat['question']})
        messages.append({"role": "assistant", "content": chat['answer']})
    
    # Build better user message
    user_message = f"""Here is content from the documents:

{context}

Question: {question}

Please answer based on the document content above."""
    
    messages.append({"role": "user", "content": user_message})
    
    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=messages,
            max_tokens=2048,
            temperature=0.1
        )
        
        answer = response.choices[0].message.content
        
        # 5. Save to history
        for file_id in file_ids:
            save_chat_history(file_id, question, answer)

        return JSONResponse(content={
            "question": question,
            "answer": answer,
            "debug_info": {
                "text_chunks_retrieved": len(text_chunks),
                "image_descriptions_retrieved": len(image_descriptions),
                "has_images": len(image_descriptions) > 0
            }
        })
        
    except Exception as e:
        error_msg = f"Groq API error: {e}"
        print(error_msg)
        return JSONResponse(content={"error": error_msg}, status_code=500)
    
@app.get("/file-ids")
def list_file_ids():
    try:
        cur = conn.cursor()
        cur.execute(f"""
            SELECT DISTINCT {TABLE_NAME}.cmetadata->>'file_id' AS file_id,
                            {TABLE_NAME}.cmetadata->>'file_name' AS file_name
            FROM {TABLE_NAME}
            ORDER BY file_name
        """)
        rows = cur.fetchall()
        cur.close()
        conn.commit()

        files = [{"file_id": r[0], "file_name": r[1]} for r in rows]
        return JSONResponse(content={"files": files}, status_code=200)

    except Exception as e:
        conn.rollback()
        return JSONResponse(content={"error": str(e)}, status_code=500)
    
@app.get("/chat-history")
def get_chat_history_endpoint(file_ids: List[str] = Query(None, alias="file_ids[]"), limit: int = 20):
    if not file_ids:
        raise HTTPException(status_code=400, detail="file_ids are required")
    
    print(f"Received file_ids: {file_ids}")  # Debug log
    
    history = get_chat_history(file_ids, limit)
    return {
        "chat_history": history,
        "file_ids": file_ids,
        "total_messages": len(history)
    }

@app.post("/clear-chat-history")
def clear_chat_history(data: dict):
    file_ids = data.get("file_ids", [])

    try:
        cur = conn.cursor()
        placeholders = ','.join(['%s'] * len(file_ids))
        cur.execute(f"DELETE FROM chat_history WHERE file_id IN ({placeholders})", file_ids)

        conn.commit()
        cur.close()

        return {
            "message": "Chat history cleared successfully",
            "file_ids": file_ids
        }
    except Exception as e:
        conn.rollback()
        return JSONResponse(content={"error": f"Failed to clear history: {str(e)}"}, status_code=500)

@app.post("/generate-flashcards")
def generate_flashcards(data: dict):
    file_ids = data.get("file_ids", [])
    num_flashcards = data.get("num_flashcards", 5)
    fill_gaps = data.get("fill_gaps", False)
    
    if not file_ids:
        raise HTTPException(status_code=400, detail="file_ids are required")
    
    # Use comprehensive search terms that work for both text and images
    search_terms = [
        "key concepts important information main topics",
        "definitions formulas equations processes",
        "diagrams charts graphs visual explanations",
        "images pictures illustrations figures",
        "content material study notes information"
    ]
    
    # Try MULTIPLE search terms to get comprehensive content
    all_text_chunks = []
    all_image_descriptions = []
    
    for term in search_terms[:3]:  # Try first 3 terms
        text_chunks, image_descriptions = retrieve_by_file_ids(file_ids, term, k=6)
        
        # Add unique text chunks
        for chunk in text_chunks:
            if chunk not in all_text_chunks:
                all_text_chunks.append(chunk)
        
        # Add unique image descriptions
        for desc in image_descriptions:
            if desc not in all_image_descriptions:
                all_image_descriptions.append(desc)
    
    print(f"DEBUG: Total text chunks collected: {len(all_text_chunks)}")
    print(f"DEBUG: Total image descriptions collected: {len(all_image_descriptions)}")
    
    # Build context with BOTH text and images
    context_parts = []
    
    if all_text_chunks:
        text_context = "TEXT CONTENT FROM NOTES/DOCUMENTS:\n" + "\n---\n".join(all_text_chunks[:10])  # Limit to 10 chunks
        context_parts.append(text_context)
    
    if all_image_descriptions:
        # Format image descriptions properly
        image_context = "VISUAL CONTENT FROM IMAGES/DIAGRAMS/FIGURES:\n"
        image_context += "\n".join([f"- {desc}" for desc in all_image_descriptions])
        context_parts.append(image_context)
    
    # If no content found at all, try one more time with empty query
    if not context_parts:
        print("WARNING: No content found with standard terms, trying empty query...")
        all_text_chunks, all_image_descriptions = retrieve_by_file_ids(file_ids, "", k=10)
        
        if all_text_chunks or all_image_descriptions:
            if all_text_chunks:
                context_parts.append("TEXT CONTENT:\n" + "\n".join(all_text_chunks))
            if all_image_descriptions:
                context_parts.append("IMAGE CONTENT:\n" + "\n".join(all_image_descriptions))
    
    context = "\n\n".join(context_parts) if context_parts else "No content found."

    if not context.strip() or context == "No content found.":
        raise HTTPException(
            status_code=400, 
            detail=f"No content found in selected files. Text chunks: {len(all_text_chunks)}, Images: {len(all_image_descriptions)}"
        )

    # Generate flashcards with the combined context
    flashcards = generate_flashcards_with_groq(context, num_flashcards, fill_gaps)

    return {
        "flashcards": flashcards,
        "file_ids_used": file_ids,
        "total_generated": len(flashcards),
        "fill_gaps_used": fill_gaps,
        "debug_info": {
            "text_chunks_used": len(all_text_chunks),
            "image_descriptions_used": len(all_image_descriptions),
            "context_length": len(context),
            "file_count": len(file_ids)
        }
    }
@app.post("/enhance-notes")
def enhance_notes(data: dict):
    """Get additional notes to fill knowledge gaps - INCLUDING IMAGES"""
    file_ids = data.get("file_ids", [])
    
    if not file_ids:
        raise HTTPException(status_code=400, detail="file_ids are required")
    
    # Get comprehensive context including images
    all_text_chunks = []
    all_image_descriptions = []
    
    # Try multiple searches
    for term in ["comprehensive understanding", "detailed information", "key concepts visual content"]:
        text_chunks, image_descriptions = retrieve_by_file_ids(file_ids, term, k=8)
        
        for chunk in text_chunks:
            if chunk not in all_text_chunks:
                all_text_chunks.append(chunk)
        
        for desc in image_descriptions:
            if desc not in all_image_descriptions:
                all_image_descriptions.append(desc)
    
    # Build context
    context_parts = []
    
    if all_text_chunks:
        text_context = "CURRENT NOTES TEXT:\n" + "\n---\n".join(all_text_chunks)
        context_parts.append(text_context)
    
    if all_image_descriptions:
        image_context = "CURRENT VISUAL CONTENT:\n" + "\n".join([f"- {desc}" for desc in all_image_descriptions])
        context_parts.append(image_context)
    
    context = "\n\n".join(context_parts) if context_parts else "No content found."
    
    if not context.strip():
        raise HTTPException(status_code=400, detail="No content found in selected files")
    
    # Generate missing notes with the full context
    result = generate_missing_notes(context)
    
    return {
        "enhanced_notes": result.get("missing_notes", []),
        "study_advice": result.get("study_advice", ""),
        "file_ids_used": file_ids,
        "original_content_summary": {
            "text_chunks": len(all_text_chunks),
            "images": len(all_image_descriptions),
            "total_content": f"{len(all_text_chunks)} text chunks, {len(all_image_descriptions)} images"
        }
    }

@app.post("/generate-mcq")
def generate_mcq(data: dict):
    file_ids = data.get("file_ids", [])
    num_questions = data.get("num_questions", 5)

    if not file_ids:
        raise HTTPException(status_code=400, detail="file_ids are required")
    
    # Use search terms that capture ALL content types
    search_terms = [
        "key concepts definitions important points",
        "formulas equations calculations data",
        "diagrams images charts visual information",
        "processes methods steps procedures",
        "facts information details explanations"
    ]
    
    # Collect content from multiple searches
    all_text_chunks = []
    all_image_descriptions = []
    
    for term in search_terms[:3]:  # Try first 3 terms
        text_chunks, image_descriptions = retrieve_by_file_ids(file_ids, term, k=6)
        
        for chunk in text_chunks:
            if chunk not in all_text_chunks:
                all_text_chunks.append(chunk)
        
        for desc in image_descriptions:
            if desc not in all_image_descriptions:
                all_image_descriptions.append(desc)
    
    print(f"MCQ DEBUG: Text chunks: {len(all_text_chunks)}, Images: {len(all_image_descriptions)}")
    
    # Build comprehensive context
    context_parts = []
    
    if all_text_chunks:
        text_context = "DOCUMENT TEXT CONTENT:\n" + "\n---\n".join(all_text_chunks[:12])  # Limit to reasonable size
        context_parts.append(text_context)
    
    if all_image_descriptions:
        image_context = "VISUAL CONTENT (Images/Diagrams/Charts):\n"
        for i, desc in enumerate(all_image_descriptions):
            # Clean up the description
            clean_desc = desc
            for prefix in ["This image shows", "The image depicts", "The image contains", "This is an image of"]:
                if desc.startswith(prefix):
                    clean_desc = desc[len(prefix):].strip()
                    break
            
            image_context += f"\nVisual {i+1}: {clean_desc}"
        context_parts.append(image_context)
    
    # Last resort if no content
    if not context_parts:
        all_text_chunks, all_image_descriptions = retrieve_by_file_ids(file_ids, "", k=15)
        
        combined = []
        if all_text_chunks:
            combined.extend(all_text_chunks)
        if all_image_descriptions:
            combined.extend(all_image_descriptions)
        
        if combined:
            context_parts.append("ALL AVAILABLE CONTENT:\n" + "\n".join(combined))
    
    context = "\n\n".join(context_parts) if context_parts else "No content found."

    if not context.strip() or context == "No content found.":
        raise HTTPException(
            status_code=400, 
            detail=f"No content available to generate questions. Found {len(all_text_chunks)} text chunks and {len(all_image_descriptions)} images."
        )

    # Generate MCQs from the combined context
    mcqs = generate_mcq_with_groq(context, num_questions)

    return {
        "mcqs": mcqs,
        "file_ids_used": file_ids,
        "total_questions": len(mcqs),
        "debug_info": {
            "text_chunks_used": len(all_text_chunks),
            "image_descriptions_used": len(all_image_descriptions),
            "context_length": len(context)
        }
    }
@app.post("/delete-file-embeddings")
def delete_file_embeddings(data: dict):
    """Delete all embeddings and chat history for given file IDs"""
    file_ids = data.get("file_ids", [])
    
    if not file_ids:
        return JSONResponse(content={"error": "file_ids are required"}, status_code=400)
    
    try:
        cur = conn.cursor()
        
        placeholders = ','.join(['%s'] * len(file_ids))
        
        deleted_embeddings = 0
        deleted_clip_embeddings = 0
        deleted_chat_history = 0
        
        cur.execute(f"DELETE FROM {TABLE_NAME} WHERE cmetadata->>'file_id' IN ({placeholders})", file_ids)
        deleted_embeddings = cur.rowcount
        
        cur.execute(f"DELETE FROM {TABLE_NAME}_clip WHERE cmetadata->>'file_id' IN ({placeholders})", file_ids)
        deleted_clip_embeddings = cur.rowcount
        
        cur.execute(f"DELETE FROM chat_history WHERE file_id IN ({placeholders})", file_ids)
        deleted_chat_history = cur.rowcount
        
        conn.commit()
        cur.close()
        
        return {
            "message": "File embeddings and chat history deleted successfully",
            "deleted_embeddings": deleted_embeddings,
            "deleted_clip_embeddings": deleted_clip_embeddings,
            "deleted_chat_history": deleted_chat_history,
            "file_ids": file_ids
        }
    except Exception as e:
        conn.rollback()
        print(f"Error deleting embeddings: {e}")
        return JSONResponse(content={"error": f"Failed to delete embeddings: {str(e)}"}, status_code=500)

@app.post("/test-image-upload")
async def test_image_upload(file: UploadFile = File(...)):
    """Test endpoint to verify image processing works"""
    filename = file.filename
    content = await file.read()
    file_stream = io.BytesIO(content)
    
    print(f"\nTEST: Uploading {filename}")
    
    # 1. Test file processing
    text, error, images, pdf_data = process_file_content(file_stream, filename)
    
    result = {
        "filename": filename,
        "file_type": filename.split('.')[-1].lower(),
        "processing_result": {
            "text_length": len(text) if text else 0,
            "images_found": len(images),
            "error": error,
            "pdf_generated": pdf_data is not None
        }
    }
    
    # 2. Test image description generation
    if images and len(images) > 0:
        print(f"Testing image description generation...")
        try:
            description = generate_image_description(images[0])
            result["image_description"] = description[:200] + "..."
            result["description_length"] = len(description)
            print(f"Generated description ({len(description)} chars)")
        except Exception as e:
            result["description_error"] = str(e)
            print(f"Description failed: {e}")
    
    # 3. Test CLIP storage
    if images and len(images) > 0:
        print(f"Testing CLIP storage...")
        try:
            clip_store = get_clip_vector_store()
            test_file_id = "test-" + str(uuid.uuid4())
            description = generate_image_description(images[0]) if 'image_description' not in result else result["image_description"]
            
            metadata = {
                "file_id": test_file_id,
                "file_name": filename,
                "image_index": 0,
                "file_type": "image",
                "content_type": "image",
                "original_content": description[:500]
            }
            
            clip_store.add_texts(
                texts=[description],
                metadatas=[metadata],
                ids=[str(uuid.uuid4())]
            )
            result["clip_storage"] = "Success"
            print(f"Stored in CLIP vector store")
            
            # Verify storage
            try:
                # Try to retrieve it
                retrieved = clip_store.similarity_search(description[:50], k=1, filter={"file_id": test_file_id})
                result["retrieval_test"] = f"Retrieved {len(retrieved)} items"
                print(f"Verified retrieval works")
            except Exception as e:
                result["retrieval_test"] = f"Retrieval failed: {e}"
                print(f"Retrieval failed: {e}")
                
        except Exception as e:
            result["clip_storage_error"] = str(e)
            print(f"CLIP storage failed: {e}")
    
    print(f"TEST COMPLETE")
    return result

@app.post("/check-database")
def check_database(data: dict):
    """Check what's in the database for given file_ids"""
    file_ids = data.get("file_ids", [])
    
    if not file_ids:
        return {"error": "No file_ids provided"}
    
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        
        placeholders = ','.join(['%s'] * len(file_ids))
        
        results = {
            "file_ids": file_ids,
            "text_store": [],
            "clip_store": []
        }
        
        # Check text store
        cur.execute(f"""
            SELECT 
                cmetadata->>'file_id' as file_id,
                cmetadata->>'file_name' as file_name,
                COUNT(*) as count,
                cmetadata->>'content_type' as content_type
            FROM {TABLE_NAME}
            WHERE cmetadata->>'file_id' IN ({placeholders})
            GROUP BY cmetadata->>'file_id', cmetadata->>'file_name', cmetadata->>'content_type'
        """, file_ids)
        
        for row in cur.fetchall():
            results["text_store"].append({
                "file_id": row[0],
                "file_name": row[1],
                "count": row[2],
                "content_type": row[3]
            })
        
        # Check CLIP store
        cur.execute(f"""
            SELECT 
                cmetadata->>'file_id' as file_id,
                cmetadata->>'file_name' as file_name,
                COUNT(*) as count,
                cmetadata->>'content_type' as content_type
            FROM {TABLE_NAME}_clip
            WHERE cmetadata->>'file_id' IN ({placeholders})
            GROUP BY cmetadata->>'file_id', cmetadata->>'file_name', cmetadata->>'content_type'
        """, file_ids)
        
        for row in cur.fetchall():
            results["clip_store"].append({
                "file_id": row[0],
                "file_name": row[1],
                "count": row[2],
                "content_type": row[3]
            })
        
        cur.close()
        conn.close()
        
        return results
        
    except Exception as e:
        return {"error": str(e)}
    


#----------------------------OCR------------------------
def extract_text_with_groq_ocr(file_stream, filename):
    """Extract text from ANY file using Groq's OCR capability"""
    client = Groq(api_key=GROQ_API_KEY)
    
    # Reset stream position
    file_stream.seek(0)
    
    # Convert file to base64
    file_bytes = file_stream.read()
    
    # For non-image files, convert to image first
    file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    if file_ext in ['pdf', 'docx', 'pptx', 'txt']:
        # Convert document to image(s) first
        return extract_text_from_document_with_ocr(file_bytes, filename, client)
    else:
        # Direct OCR for image files
        return extract_text_from_image(file_bytes, client)


def extract_text_from_image(image_bytes, client):
    """Extract text from image bytes using Groq"""
    img_b64 = base64.b64encode(image_bytes).decode('utf-8')
    
    # Simple OCR prompt - JUST EXTRACT TEXT
    response = client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": "Extract ALL text from this image. Return ONLY the text, no explanations."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}}
            ]
        }],
        max_tokens=2000,
        temperature=0.1
    )
    
    return response.choices[0].message.content.strip()


def extract_text_from_document_with_ocr(file_bytes, filename, client):
    """Extract text from documents by converting pages to images and OCRing each"""
    # Convert PDF/DOC to images
    images = []
    
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        # Convert PDF pages to images
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Higher resolution for OCR
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            images.append(img)
        
        doc.close()
    
    elif file_ext in ['docx', 'pptx']:
        # Convert office docs to PDF then to images
        # (Use your existing convert_office_to_pdf function)
        pdf_data, error = convert_office_to_pdf(io.BytesIO(file_bytes), filename)
        if pdf_data:
            # Convert PDF to images as above
            doc = fitz.open(stream=pdf_data, filetype="pdf")
            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                images.append(img)
            doc.close()
    
    # OCR each image
    all_text = []
    for i, img in enumerate(images):
        print(f"  OCRing page {i+1}/{len(images)}...")
        
        # Convert image to bytes
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_b64 = base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')
        
        # Extract text
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract ALL text from this document page. Return ONLY the text, no explanations."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}}
                ]
            }],
            max_tokens=1000,
            temperature=0.1
        )
        
        page_text = response.choices[0].message.content.strip()
        if page_text:
            all_text.append(f"--- Page {i+1} ---\n{page_text}")
    
    return "\n\n".join(all_text)

# ---------------- RUN ----------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))  # fallback for local
    uvicorn.run(app, host="0.0.0.0", port=port)


