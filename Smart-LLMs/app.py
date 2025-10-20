from flask import Flask, request, jsonify
import os
import PyPDF2
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from transformers import pipeline
import numpy as np
import faiss
import re

app = Flask(__name__)

class RealChatbot:
    def __init__(self):
        print("Starting Chatbot!")
        
        # Load embedding model for search
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Load a model that can read and explain content
        self.explainer = pipeline(
            "text2text-generation",
            model="google/flan-t5-base",  # Good at following instructions
            max_length=500
        )
        
        # Initialize storage
        self.documents = []
        self.full_text = ""
        self.index = None

        self.current_file = None
        
    def extract_text(self, file_path, filename):
        """Extract text from files"""
        text = ""
        
        try:
            if filename.endswith('.pdf'):
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                        
            elif filename.endswith('.docx'):
                doc = Document(file_path)
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                    
            elif filename.endswith('.pptx'):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
            print(f"Extracted {len(text)} characters from {filename}")
            self.full_text = text
            return text
            
        except Exception as e:
            print(f"Error: {e}")
            return ""
    
    def chunk_text(self, text):
        """Split text into meaningful chunks"""
        # Split by paragraphs or sentences
        paragraphs = text.split('\n\n')
        chunks = []
        
        for para in paragraphs:
            if len(para.strip()) > 50:  # Only keep substantial paragraphs
                chunks.append(para.strip())
        
        return chunks
    
    def process_file(self, file_path, filename):
        """Process uploaded file"""
        text = self.extract_text(file_path, filename)
        if not text:
            return False
        
        # Create chunks
        chunks = self.chunk_text(text)
        self.documents = chunks
        self.current_file = filename
        
        # Create embeddings for search
        embeddings = self.embedding_model.encode(chunks)
        
        # Create FAISS index
        dimension = embeddings.shape[1]
        self.index = faiss.IndexFlatIP(dimension)
        faiss.normalize_L2(embeddings)
        self.index.add(embeddings.astype(np.float32))
        
        print(f"Processed {len(chunks)} chunks from {filename}")
        return True
    
    def find_relevant_text(self, question, top_k=3):
        """Find most relevant text for the question"""
        if not self.index:
            return []
        
        # Encode question
        question_embedding = self.embedding_model.encode([question])
        faiss.normalize_L2(question_embedding)
        
        # Search
        scores, indices = self.index.search(question_embedding.astype(np.float32), top_k)
        
        results = []
        for idx, score in zip(indices[0], scores[0]):
            if idx < len(self.documents) and score > 0.2:
                results.append(self.documents[idx])
        
        return results
    
    def generate_real_explanation(self, question, relevant_texts):
        """Use the model to read and explain the content"""
        
        if not relevant_texts:
            return "I couldn't find specific information about this in your document. Try asking about topics that might be in your notes."
        
        # Combine relevant texts as context
        context = "\n".join(relevant_texts[:3])
        
        # Create a prompt that makes the model read and explain
        prompt = f"""
Based on the following text from a document, please explain the answer to this question: "{question}"

Document text:
{context}

Please provide a clear explanation using only the information from the document above. Be helpful and educational.
Explanation:
"""
        
        try:
            # Let the model actually process and explain
            response = self.explainer(
                prompt,
                max_length=400,
                do_sample=True,
                temperature=0.7
            )
            
            explanation = response[0]['generated_text']
            
            # If the explanation is too short or generic, provide a fallback
            if len(explanation.strip()) < 50:
                return self._create_fallback_explanation(question, relevant_texts)
            
            return explanation
            
        except Exception as e:
            print(f"Model error: {e}")
            return self._create_fallback_explanation(question, relevant_texts)
    
    def _create_fallback_explanation(self, question, relevant_texts):
        """Fallback explanation using actual document content"""
        # Use the actual document content to create explanation
        context = "\n".join([f"• {text}" for text in relevant_texts[:2]])
        
        return f"""Based on your document, here's what I found:

**Relevant content from your notes:**
{context}

**Explanation:**
The document provides information about {question.lower()}. The key points shown above are the most relevant sections that address your question.

To understand this better, focus on the specific details mentioned in these sections of your document."""

# Initialize chatbot
chatbot = RealChatbot()

@app.route('/upload', methods=['POST'])
def upload_file():
    """Upload a file"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Save and process file
        file_path = f"temp_{file.filename}"
        file.save(file_path)
        
        success = chatbot.process_file(file_path, file.filename)
        os.remove(file_path)
        
        if not success:
            return jsonify({'error': 'Could not process file'}), 400
        
        return jsonify({
            'success': True,
            'message': 'File uploaded successfully! Ready for real explanations.',
            'filename': file.filename,
            'chunks': len(chatbot.documents)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    """Ask a question - model will ACTUALLY read your document"""
    try:
        data = request.json
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Please provide a question'}), 400
        
        if not chatbot.documents:
            return jsonify({'error': 'Please upload a document first'}), 400
        
        # Find relevant text using vector search
        relevant_texts = chatbot.find_relevant_text(question)
        
        # Generate REAL explanation using the model
        answer = chatbot.generate_real_explanation(question, relevant_texts)
        
        return jsonify({
            'question': question,
            'answer': answer,
            'found_references': len(relevant_texts),
            'explanation_type': 'model_generated'
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/summarize', methods=['POST'])
def summarize_document():
    """Get a real summary of the document"""
    try:
        if not chatbot.full_text:
            return jsonify({'error': 'No document uploaded'}), 400
        
        # Use the model to generate a real summary
        prompt = f"Please summarize the following document in a clear, educational way:\n\n{chatbot.full_text[:2000]}\n\nSummary:"
        
        response = chatbot.explainer(
            prompt,
            max_length=300,
            do_sample=True,
            temperature=0.7
        )
        
        summary = response[0]['generated_text']
        
        return jsonify({
            'summary': summary,
            'source': chatbot.current_file
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({
        'status': 'ready',
        'has_document': chatbot.current_file is not None,
        'current_file': chatbot.current_file,
        'model': 'google/flan-t5-base'
    })

if __name__ == '__main__':
    print(" Chatbot Ready!")
    print("\nHow to use:")
    print("1. POST /upload - Upload a PDF/DOCX/PPTX file")
    print("2. POST /ask - Ask questions (AI reads your document)")
    print("3. POST /summarize - Get AI-generated summary")
    print("4. GET /health - Check status")
    app.run(debug=True, host='0.0.0.0', port=5000)