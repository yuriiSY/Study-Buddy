from flask import Flask, request, jsonify
from dotenv import load_dotenv
import os
import PyPDF2
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import re
import google.generativeai as genai


app = Flask(__name__)

# Configure Gemini - REPLACE WITH YOUR API KEY
GEMINI_API_KEY = os.getenv('API_KEY')
genai.configure(api_key=GEMINI_API_KEY)

class SmartHybridAssistant:
    def __init__(self):
        print("Starting Smart Assistant")
        
        # Load embedding model for document search
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Initialize Gemini model
        try:
            self.gemini_model = genai.GenerativeModel('gemini-2.0-flash-exp')
            print("Gemini model loaded successfully!")
        except Exception as e:
            print(f"Error loading Gemini: {e}")
            raise
        
        # Initialize storage
        self.documents = []
        self.full_text = ""
        self.index = None
        self.current_file = None
        
    def extract_text(self, file_path, filename):
        """Extract text from files"""
        text = ""
        
        try:
            if filename.endswith('.pdf'):
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                        
            elif filename.endswith('.docx'):
                doc = Document(file_path)
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                    
            elif filename.endswith('.pptx'):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
            print(f"Extracted {len(text)} characters from {filename}")
            self.full_text = text
            return text
            
        except Exception as e:
            print(f"Error: {e}")
            return ""
    
    def chunk_text(self, text):
        """Split text into meaningful chunks"""
        paragraphs = text.split('\n\n')
        chunks = []
        
        for para in paragraphs:
            if len(para.strip()) > 50:
                chunks.append(para.strip())
        
        return chunks
    
    def process_file(self, file_path, filename):
        """Process uploaded file"""
        text = self.extract_text(file_path, filename)
        if not text:
            return False
        
        chunks = self.chunk_text(text)
        self.documents = chunks
        self.current_file = filename
        
        # Create embeddings for search
        embeddings = self.embedding_model.encode(chunks)
        
        # Create FAISS index
        dimension = embeddings.shape[1]
        self.index = faiss.IndexFlatIP(dimension)
        faiss.normalize_L2(embeddings)
        self.index.add(embeddings.astype(np.float32))
        
        print(f"Processed {len(chunks)} chunks from {filename}")
        return True
    
    def find_relevant_text(self, question, top_k=5):
        """Find most relevant text for the question"""
        if not self.index:
            return []
        
        question_embedding = self.embedding_model.encode([question])
        faiss.normalize_L2(question_embedding)
        
        scores, indices = self.index.search(question_embedding.astype(np.float32), top_k)
        
        results = []
        for idx, score in zip(indices[0], scores[0]):
            if idx < len(self.documents) and score > 0.1:
                results.append({
                    'text': self.documents[idx],
                    'score': float(score)
                })
        
        # Sort by relevance score
        results.sort(key=lambda x: x['score'], reverse=True)
        return results
    
    def analyze_sufficiency(self, question, relevant_texts):
        """Analyze if the notes have sufficient information"""
        if not relevant_texts:
            return "insufficient", 0
        
        # Calculate average relevance score
        avg_score = sum([item['score'] for item in relevant_texts]) / len(relevant_texts)
        
        # Check if we have high-quality matches
        high_quality_matches = [item for item in relevant_texts if item['score'] > 0.3]
        
        if len(high_quality_matches) >= 2 and avg_score > 0.25:
            return "sufficient", avg_score
        elif len(high_quality_matches) >= 1:
            return "partial", avg_score
        else:
            return "insufficient", avg_score
    
    def generate_smart_answer(self, question, relevant_texts):
        """Smart hybrid approach: notes first, supplement when needed"""
        
        # Analyze if notes are sufficient
        sufficiency, avg_score = self.analyze_sufficiency(question, relevant_texts)
        
        # Extract just the text from relevant items
        context_texts = [item['text'] for item in relevant_texts]
        context = "\n\n".join(context_texts[:3])  # Use top 3 most relevant
        
        print(f"Notes sufficiency: {sufficiency} (avg score: {avg_score:.3f})")
        
        if sufficiency == "sufficient":
            # Primary: Use notes with minimal supplementation
            prompt = f"""Answer this question using the provided notes, and supplement with basic general knowledge if needed to connect ideas.

QUESTION: {question}

STUDENT'S NOTES:
{context}

Please provide an answer that:
1. FIRST uses the information from the notes above
2. Adds  general knowledge as needed to explain in more easy and practical way
3. Stays focused on the topic

Answer:"""
        
        elif sufficiency == "partial":
            # Balanced: Use notes as base, supplement more
            prompt = f"""Answer this question using the provided notes as a starting point, and supplement with general knowledge to provide a complete explanation.

QUESTION: {question}

RELEVANT NOTES CONTENT:
{context}

Please provide an answer that:
1. Starts with what the notes cover
2. Supplements with general knowledge to fill gaps
3. Explains how the general knowledge connects to the notes
4. Provides a comprehensive understanding

Answer:"""
        
        else:  # insufficient
            # Secondary: Use general knowledge but relate to notes context
            if context:
                prompt = f"""The notes don't have sufficient information about this specific question. Using general knowledge, provide an answer that might be relevant to the overall context of the notes.

QUESTION: {question}

CONTEXT FROM NOTES (for reference):
{context}

Based on general knowledge, provide an answer that could be relevant to this student's studies:"""
            else:
                prompt = f"""The notes don't contain information about this topic. Provide a helpful educational answer using general knowledge.

QUESTION: {question}

Answer as a helpful tutor:"""
        
        try:
            response = self.gemini_model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.3,
                    "max_output_tokens": 600,
                }
            )
            
            answer = response.text
            
            # Add metadata about source
            if sufficiency == "sufficient":
                source_info = "Primarily from your notes"
            elif sufficiency == "partial":
                source_info = "Based on your notes + general knowledge"
            else:
                source_info = "General knowledge (notes don't cover this)"
            
            return {
                'answer': answer,
                'source': source_info,
                'sufficiency': sufficiency,
                'relevance_score': avg_score,
                'found_chunks': len(relevant_texts)
            }
            
        except Exception as e:
            print(f"Gemini error: {e}")
            # Fallback: just return the notes content
            if context_texts:
                return {
                    'answer': "Based on your notes:\n\n" + "\n\n".join(context_texts[:3]),
                    'source': 'Direct from notes (fallback)',
                    'sufficiency': 'partial',
                    'relevance_score': avg_score,
                    'found_chunks': len(relevant_texts)
                }
            else:
                return {
                    'answer': "No relevant information found in your notes.",
                    'source': 'No notes content',
                    'sufficiency': 'insufficient',
                    'relevance_score': 0,
                    'found_chunks': 0
                }
    
    def explain_with_context(self, concept):
        """Explain a concept using notes context + supplementation"""
        relevant_texts = self.find_relevant_text(concept, top_k=5)
        
        # Extract context
        context_texts = [item['text'] for item in relevant_texts]
        context = "\n\n".join(context_texts[:3])
        
        prompt = f"""Explain this concept in a way that connects to the student's study materials.

CONCEPT: {concept}

CONTEXT FROM STUDENT'S NOTES:
{context}

Please provide an explanation that:
1. First addresses what the notes cover about this concept
2. Supplements with general knowledge to provide a complete understanding
3. Helps the student to understand in a practical and easy way

Explanation:"""
        
        try:
            response = self.gemini_model.generate_content(prompt)
            return response.text
        except Exception as e:
            return f"Error generating explanation: {e}"

# Initialize chatbot
chatbot = SmartHybridAssistant()

@app.route('/upload', methods=['POST'])
def upload_file():
    """Upload a file"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        file_path = f"temp_{file.filename}"
        file.save(file_path)
        
        success = chatbot.process_file(file_path, file.filename)
        os.remove(file_path)
        
        if not success:
            return jsonify({'error': 'Could not process file'}), 400
        
        return jsonify({
            'success': True,
            'message': 'File uploaded successfully! I will use your notes first, then supplement intelligently.',
            'filename': file.filename,
            'chunks': len(chatbot.documents)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    """Ask questions - smart hybrid approach"""
    try:
        data = request.json
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Please provide a question'}), 400
        
        if not chatbot.documents:
            return jsonify({'error': 'Please upload a document first'}), 400
        
        # Find relevant text
        relevant_texts = chatbot.find_relevant_text(question)
        
        # Generate smart hybrid answer
        result = chatbot.generate_smart_answer(question, relevant_texts)
        
        return jsonify({
            'question': question,
            'answer': result['answer'],
            'source': result['source'],
            'notes_sufficiency': result['sufficiency'],
            'relevance_score': round(result['relevance_score'], 3),
            'chunks_found': result['found_chunks']
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/explain', methods=['POST'])
def explain_concept():
    """Explain a concept with notes context"""
    try:
        data = request.json
        concept = data.get('concept', '').strip()
        
        if not concept:
            return jsonify({'error': 'Please provide a concept'}), 400
        
        if not chatbot.documents:
            return jsonify({'error': 'Please upload a document first'}), 400
        
        explanation = chatbot.explain_with_context(concept)
        
        return jsonify({
            'concept': concept,
            'explanation': explanation,
            'approach': 'notes_context_plus_general'
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/analyze_notes', methods=['GET'])
def analyze_notes():
    """Analyze what topics are covered in the notes"""
    try:
        if not chatbot.documents:
            return jsonify({'error': 'No notes available'}), 400
        
        # Sample some content to analyze
        sample_text = "\n".join(chatbot.documents[:10])
        
        prompt = f"""Based on the following notes content, what are the main topics and subjects covered?

NOTES CONTENT:
{sample_text}

Please list the main topics and subjects:"""
        
        try:
            response = chatbot.gemini_model.generate_content(prompt)
            analysis = response.text
        except:
            analysis = "Analysis unavailable"
        
        return jsonify({
            'filename': chatbot.current_file,
            'total_chunks': len(chatbot.documents),
            'topics_analysis': analysis,
            'sample_content': sample_text[:500] + "..." if len(sample_text) > 500 else sample_text
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    print("\nEndpoints:")
    print("   POST /upload        - Upload study materials")
    print("   POST /ask           - Ask questions (smart hybrid approach)")
    print("   POST /explain       - Explain concepts with notes context")
    print("   GET  /analyze_notes - See what topics your notes cover")
    
    print("\nHow it works:")
    print("   • Analyzes if your notes have sufficient information")
    print("   • Uses notes content as primary source")
    print("   • Intelligently supplements when notes are insufficient")
    print("   • Connects general knowledge to your notes context")
    
    app.run(debug=True, host='0.0.0.0', port=5000)